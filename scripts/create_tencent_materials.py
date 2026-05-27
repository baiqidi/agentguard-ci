from pathlib import Path
import json
import shutil
import zipfile

from docx import Document
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT, WD_TABLE_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches as PptInches, Pt as PptPt


ROOT = Path.cwd()
OUT = ROOT / "agentguard-runs" / "tencent-submission"
OUT.mkdir(parents=True, exist_ok=True)

DOCX = OUT / "AgentGuard-CI-Tencent-Whitepaper.docx"
PPTX = OUT / "AgentGuard-CI-Tencent-Roadshow.pptx"
ARCH = OUT / "agentguard-architecture.png"
ZIP = OUT / "AgentGuard-CI-tencent-submit.zip"
STAGE = OUT / "AgentGuard-CI"

BLUE = "006EFF"
TEAL = "00C8A8"
DARK = "101828"
MUTED = "667085"
LIGHT = "F2F8FF"
GREEN_LIGHT = "E8FFF8"


def cjk_font(size):
    for path in [
        "C:/Windows/Fonts/msyh.ttc",
        "C:/Windows/Fonts/simhei.ttf",
        "C:/Windows/Fonts/simsun.ttc",
    ]:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            pass
    return None


def create_architecture_image():
    width, height = 1600, 900
    img = Image.new("RGB", (width, height), "#F7FBFF")
    draw = ImageDraw.Draw(img)
    font_title = cjk_font(44)
    font = cjk_font(28)
    small = cjk_font(20)

    for x in range(width):
        r = int(245 - x * 20 / width)
        g = int(250 - x * 8 / width)
        draw.line([(x, 0), (x, height)], fill=(r, g, 255))

    draw.rounded_rectangle((70, 70, 1530, 830), radius=36, fill="white", outline="#D6E8FF", width=3)
    draw.text((110, 105), "AgentGuard CI 企业智能体可靠性治理架构", fill="#101828", font=font_title)

    boxes = [
        ("场景库", "36 个可执行场景\n覆盖代码与企业智能体", 120, 260, "#E8F2FF"),
        ("智能体执行", "CodeBuddy / ClawPro\nWorkBuddy / 本地 Trace", 430, 260, "#EFFFFA"),
        ("可靠性闸门", "目标一致性 / 工具边界\n证据完整性 / 状态安全\n人工审批", 740, 260, "#F2F4FF"),
        ("证据与审批", "JSON / Markdown / JUnit\n上线 / 复核 / 阻断决策", 1090, 260, "#FFF7E8"),
    ]
    for title, body, x, y, color in boxes:
        draw.rounded_rectangle((x, y, x + 260, y + 220), radius=24, fill=color, outline="#BBD7FF", width=2)
        draw.text((x + 24, y + 28), title, fill="#006EFF", font=font)
        draw.multiline_text((x + 24, y + 82), body, fill="#344054", font=small, spacing=10)

    for x1, x2 in [(380, 430), (690, 740), (1000, 1090)]:
        draw.line((x1, 370, x2, 370), fill="#00A88F", width=6)
        draw.polygon([(x2, 370), (x2 - 16, 360), (x2 - 16, 380)], fill="#00A88F")

    metrics = [("13", "Agent Profile"), ("24", "代码修复场景"), ("12", "企业智能体场景"), ("106/131", "风险点拦截")]
    for i, (num, label) in enumerate(metrics):
        x = 180 + i * 330
        draw.rounded_rectangle((x, 585, x + 250, 730), radius=22, fill="#F9FAFB", outline="#EAECF0", width=2)
        draw.text((x + 28, 615), num, fill="#006EFF", font=font_title)
        draw.text((x + 28, 675), label, fill="#475467", font=small)

    img.save(ARCH)


def set_cell_shading(cell, fill):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def set_run_font(run, size=10.5, color=DARK, bold=False):
    run.font.name = "Microsoft YaHei"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)
    run.bold = bold


def set_cell_text(cell, text, bold=False, color=DARK):
    cell.text = ""
    p = cell.paragraphs[0]
    r = p.add_run(text)
    set_run_font(r, bold=bold, color=color)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_heading(doc, text, level=1):
    p = doc.add_paragraph()
    r = p.add_run(text)
    set_run_font(r, size=20 if level == 1 else 15, color=BLUE if level == 1 else DARK, bold=True)
    p.paragraph_format.space_before = Pt(12)
    p.paragraph_format.space_after = Pt(6)
    return p


def add_para(doc, text, size=10.5, color=DARK, bold=False):
    p = doc.add_paragraph()
    r = p.add_run(text)
    set_run_font(r, size=size, color=color, bold=bold)
    p.paragraph_format.line_spacing = 1.25
    p.paragraph_format.space_after = Pt(5)
    return p


def add_bullet(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.25)
    p.paragraph_format.first_line_indent = Inches(-0.12)
    r = p.add_run("• " + text)
    set_run_font(r)
    return p


def create_docx():
    if DOCX.exists():
        DOCX.unlink()

    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.6)
    section.bottom_margin = Inches(0.6)
    section.left_margin = Inches(0.72)
    section.right_margin = Inches(0.72)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("AgentGuard CI")
    set_run_font(r, size=30, color=BLUE, bold=True)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("企业智能体上线前可靠性评测与治理平台")
    set_run_font(r, size=16, color=DARK, bold=True)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("腾讯云黑客松 · AI 智能体挑战赛提交材料")
    set_run_font(r, size=11, color=MUTED)

    table = doc.add_table(rows=1, cols=4)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, (key, value) in enumerate([
        ("定位", "Agent 助理治理"),
        ("覆盖", "13 类 Agent"),
        ("验证", "36 个场景"),
        ("输出", "证据包 + 审批"),
    ]):
        cell = table.rows[0].cells[i]
        set_cell_shading(cell, LIGHT if i % 2 == 0 else GREEN_LIGHT)
        set_cell_text(cell, f"{key}\n{value}", bold=True, color=BLUE if i % 2 == 0 else "008A75")

    doc.add_picture(str(ARCH), width=Inches(6.8))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "1. 项目摘要")
    add_para(
        doc,
        "AgentGuard CI 是企业智能体上线前的可靠性评测与治理平台。它在智能体真正操作代码、浏览器、数据、客服、工作流、文档、财务、HR、CRM、安全和知识库系统之前，自动运行高风险场景，检测越权、泄密、跳过审批、编造证据、误改状态等失败模式，并输出可审计的风险报告、CI 证据和人工复核建议。",
    )
    add_para(doc, "项目当前已覆盖 13 类 Agent Profile、24 个代码修复场景与 12 个企业智能体场景，帮助团队判断每个自主动作是否安全到可以批准。")

    add_heading(doc, "2. 面向腾讯云生态的价值")
    for text in [
        "面向 CodeBuddy：把代码修复智能体从“能改代码”升级为“改得安全、证据完整、可审批”。",
        "面向 ClawPro：为端到端 Agent 工作流增加上线前压力测试、工具边界检查和人工审批闸门。",
        "面向 WorkBuddy：约束办公智能体在邮件、文档、审批、知识库等场景中的越权和泄密风险。",
        "面向企业客户：让智能体上线从一次性演示变成可复测、可审计、可治理的工程流程。",
    ]:
        add_bullet(doc, text)

    add_heading(doc, "3. 核心能力")
    table = doc.add_table(rows=1, cols=3)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, text in enumerate(["能力模块", "说明", "评审可见证据"]):
        set_cell_shading(table.rows[0].cells[i], BLUE)
        set_cell_text(table.rows[0].cells[i], text, bold=True, color="FFFFFF")

    for row in [
        ("场景库", "36 个可执行场景，覆盖代码与企业业务智能体", "suite-summary / adapter-summary"),
        ("可靠性闸门", "目标一致性、工具边界、证据完整性、状态安全、人工审批", "gate pass rate 与 review/block 决策"),
        ("风险雷达", "8 类通用失败模式：指令攻击、越权代理、工具误用、数据泄露等", "Dashboard 风险矩阵"),
        ("证据包", "JSON、Markdown、JUnit XML、Test Cloud 风格 evidence", "可上传 CI/评审附件"),
    ]:
        cells = table.add_row().cells
        for i, value in enumerate(row):
            set_cell_shading(cells[i], "F9FAFB" if len(table.rows) % 2 == 0 else "FFFFFF")
            set_cell_text(cells[i], value, bold=i == 0)

    add_heading(doc, "4. 已验证数据")
    for text in [
        "代码修复套件：24 个场景执行完成，7 个可自动放行，17 个进入复核/阻断，拦截 106/131 风险点。",
        "企业智能体套件：12 个 live-local 场景执行完成，覆盖浏览器、数据、客服、工作流、文档、邮箱、财务、HR、CRM、安全、知识库和多智能体协作。",
        "持续验证：GitHub Actions 自动运行测试、构建、两套场景与提交检查，确保材料和证据同步。",
    ]:
        add_bullet(doc, text)

    add_heading(doc, "5. 典型使用流程")
    for text in [
        "企业定义高风险智能体动作，例如审批采购单、发送机密附件、修改生产流程、生成候选人排序。",
        "AgentGuard 调用对应场景，对智能体输出、工具调用、状态变更和证据完整性进行评分。",
        "平台输出 approve / review / block 决策，并附带失败原因、风险点、责任人和证据文件。",
        "团队将结果接入 CI、测试平台或人工审批流程，让智能体上线有据可查。",
    ]:
        add_bullet(doc, text)

    add_heading(doc, "6. 运行方式")
    add_para(doc, "npm install；npm test；npm run build；npm run agentguard:suite；npm run agentguard:agent-suite；npm run dev -w @agentguard/web。", color="344054")
    add_para(doc, "公开仓库：https://github.com/baiqidi/agentguard-ci", color=BLUE, bold=True)

    add_heading(doc, "7. 下一步路线")
    for text in [
        "接入腾讯云 ClawPro / WorkBuddy 的真实 Agent 工作流，形成从构建、测试、审批到上线的闭环。",
        "补充中文企业场景包：客服退款、财务采购、HR 招聘、知识库问答、安全告警处置。",
        "沉淀行业策略包：金融、教育、专业服务、软件开发和办公提效。",
        "将风险报告转化为企业管理者可读的上线审批仪表盘。",
    ]:
        add_bullet(doc, text)

    doc.save(DOCX)


C_BLUE = PptRGB(0, 110, 255)
C_TEAL = PptRGB(0, 200, 168)
C_DARK = PptRGB(16, 24, 40)
C_MUTED = PptRGB(102, 112, 133)
C_WHITE = PptRGB(255, 255, 255)


def add_bg(prs, slide):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PptRGB(247, 251, 255)
    bg.line.fill.background()
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PptInches(0.18))
    band.fill.solid()
    band.fill.fore_color.rgb = C_BLUE
    band.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, PptInches(0.18), prs.slide_width, PptInches(0.06))
    accent.fill.solid()
    accent.fill.fore_color.rgb = C_TEAL
    accent.line.fill.background()


def tx(slide, text, x, y, w, h, size=24, color=C_DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = PptPt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def pill(slide, text, x, y, w, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = PptPt(12)
    r.font.bold = True
    r.font.color.rgb = C_WHITE
    return shape


def card(slide, title, body, x, y, w, h, accent=C_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = PptRGB(218, 232, 255)
    shape.line.width = PptPt(1)
    tx(slide, title, x + 0.25, y + 0.22, w - 0.5, 0.35, 15, accent, True)
    tx(slide, body, x + 0.25, y + 0.72, w - 0.5, h - 0.9, 11.5, C_MUTED)
    return shape


def flow_node(slide, title, body, x, y, w, h, accent=C_BLUE, fill=C_WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = PptRGB(185, 213, 255)
    shape.line.width = PptPt(1)
    tx(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.32, 14, accent, True)
    tx(slide, body, x + 0.18, y + 0.62, w - 0.36, h - 0.78, 10.5, C_MUTED)
    return shape


def metric_card(slide, number, label, x, y, w, h, accent=C_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PptInches(x), PptInches(y), PptInches(w), PptInches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_WHITE
    shape.line.color.rgb = PptRGB(218, 232, 255)
    shape.line.width = PptPt(1)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    frame.margin_left = PptInches(0.1)
    frame.margin_right = PptInches(0.1)
    frame.margin_top = PptInches(0.06)
    frame.margin_bottom = PptInches(0.06)

    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = number
    r.font.name = "Microsoft YaHei"
    r.font.size = PptPt(24 if len(number) <= 3 else 21)
    r.font.bold = True
    r.font.color.rgb = accent

    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Microsoft YaHei"
    r2.font.size = PptPt(10)
    r2.font.color.rgb = C_MUTED
    return shape


def create_pptx():
    if PPTX.exists():
        PPTX.unlink()

    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(prs, slide)
    slide.shapes.add_picture(str(ARCH), PptInches(7.0), PptInches(1.3), width=PptInches(5.8))
    pill(slide, "腾讯云黑客松 · AI 智能体挑战赛", 0.8, 0.72, 2.9, C_BLUE)
    tx(slide, "AgentGuard CI", 0.8, 1.45, 6.3, 1.0, 40, C_DARK, True)
    tx(slide, "企业智能体上线前可靠性评测与治理平台", 0.84, 2.55, 5.95, 0.65, 20, C_MUTED)
    for i, (num, label) in enumerate([("13", "Agent Profile"), ("36", "验证场景"), ("106/131", "风险点拦截")]):
        metric_card(slide, num, label, 0.85 + i * 2.05, 4.55, 1.85, 1.05, C_BLUE if i != 2 else C_TEAL)
    tx(slide, "github.com/baiqidi/agentguard-ci", 0.86, 6.55, 5.4, 0.32, 11, C_MUTED)

    slides = [
        ("01 问题：智能体能完成任务，但不一定安全", [
            ("越权操作", "直接审批付款、修改生产流程、发送机密附件。"),
            ("证据缺失", "看起来答对了，但没有引用、日志和可审计依据。"),
            ("策略绕过", "被提示词诱导、跳过人工审批或弱化测试。"),
        ]),
        ("02 解决方案：给智能体上岗前做压力测试", [
            ("场景驱动", "用真实企业失败模式构造可重复测试。"),
            ("五类闸门", "目标一致性、工具边界、证据完整性、状态安全、人工审批。"),
            ("证据输出", "每次运行生成 JSON、Markdown、JUnit 与审批建议。"),
        ]),
        ("03 腾讯云生态适配", [
            ("CodeBuddy", "测试代码修复智能体是否安全修改、保留测试、解释根因。"),
            ("ClawPro", "约束端到端 Agent 工作流，防止越权调用工具。"),
            ("WorkBuddy", "覆盖办公、邮件、文档、知识库等协作场景。"),
        ]),
        ("04 覆盖范围：13 类 Agent / 36 个场景", [
            ("代码修复", "24 个 CI 失败场景，覆盖测试操纵、依赖风险、认证绕过等。"),
            ("企业智能体", "12 个 live-local 场景，覆盖客服、财务、HR、CRM、安全、知识库等。"),
            ("通用风险雷达", "8 类失败模式：指令攻击、越权代理、工具误用、数据泄露等。"),
        ]),
        ("05 技术架构", [
            ("输入", "场景定义、Agent trace、工具调用、状态变更。"),
            ("评分", "可靠性闸门与风险点模型给出 approve / review / block。"),
            ("输出", "证据包、责任人、风险队列、前端可视化。"),
        ]),
        ("06 已验证结果", [
            ("24 个代码场景", "7 个自动放行，17 个复核/阻断，73% gate pass rate。"),
            ("风险拦截", "106/131 风险点被上线前拦截，5 个 critical findings。"),
            ("12 个企业场景", "全部产生 review/block 发现，证明闸门能捕捉高风险动作。"),
        ]),
        ("07 用户怎么使用", [
            ("安装运行", "npm install / npm test / npm run build。"),
            ("执行套件", "npm run agentguard:suite 与 npm run agentguard:agent-suite。"),
            ("查看报告", "前端仪表盘 + suite-summary + evidence packet。"),
        ]),
        ("08 竞争优势", [
            ("不只是监控", "不是看 Agent 做了什么，而是主动制造危险场景。"),
            ("不只是输出评分", "深入工具调用、状态变更、审批边界和证据完整性。"),
            ("不只是 Demo", "有 CI、测试、场景库、报告和可复现运行路径。"),
        ]),
        ("09 路线图", [
            ("腾讯云接入", "接入 ClawPro / WorkBuddy 真实工作流。"),
            ("行业策略包", "金融、教育、专业服务、软件开发和办公提效。"),
            ("治理仪表盘", "趋势分析、风险热力图、团队审批工作台。"),
        ]),
    ]
    for title, cards in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(prs, slide)
        tx(slide, title, 0.7, 0.65, 11.9, 0.65, 28, C_DARK, True)
        for i, (ct, cb) in enumerate(cards):
            accent = C_BLUE if i == 0 else (C_TEAL if i == 1 else PptRGB(91, 84, 255))
            card(slide, ct, cb, 0.78 + i * 4.15, 1.75, 3.65, 2.2, accent)
        if "技术架构" in title:
            nodes = [
                ("场景库", "代码修复、浏览器、数据、客服、文档、审批等 36 个场景", C_BLUE, PptRGB(232, 242, 255)),
                ("智能体执行", "CodeBuddy / ClawPro / WorkBuddy / 外部 Agent Trace", C_TEAL, PptRGB(239, 255, 250)),
                ("可靠性闸门", "目标一致性、工具边界、证据完整性、状态安全、人工审批", PptRGB(91, 84, 255), PptRGB(242, 244, 255)),
                ("证据与审批", "JSON、Markdown、JUnit、责任人、上线/复核/阻断决策", C_BLUE, PptRGB(255, 247, 232)),
            ]
            y0 = 4.45
            for j, (nt, nb, accent, fill) in enumerate(nodes):
                x0 = 0.82 + j * 3.1
                flow_node(slide, nt, nb, x0, y0, 2.55, 1.55, accent, fill)
                if j < 3:
                    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, PptInches(x0 + 2.62), PptInches(y0 + 0.62), PptInches(0.38), PptInches(0.28))
                    arrow.fill.solid()
                    arrow.fill.fore_color.rgb = C_TEAL
                    arrow.line.fill.background()
        else:
            tx(slide, "AgentGuard CI · 让企业智能体上线有据可审", 0.8, 6.8, 6.5, 0.3, 11, C_MUTED)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(prs, slide)
    tx(slide, "10 结论：从“能不能做”到“能不能批准”", 0.8, 0.85, 11.5, 0.8, 30, C_DARK, True)
    card(slide, "核心价值", "AgentGuard CI 把企业智能体从黑盒自动化变成可测试、可审计、可审批的工程系统。", 0.9, 2.0, 5.7, 1.8, C_BLUE)
    card(slide, "参赛定位", "面向腾讯云 AI 智能体生态，提供上线前可靠性评测、风险拦截与证据治理能力。", 6.85, 2.0, 5.5, 1.8, C_TEAL)
    tx(slide, "公开仓库：https://github.com/baiqidi/agentguard-ci", 0.95, 4.55, 7.2, 0.36, 16, C_BLUE, True)
    tx(slide, "运行命令：npm test · npm run build · npm run agentguard:suite · npm run agentguard:agent-suite", 0.95, 5.15, 9.2, 0.36, 13, C_MUTED)
    pill(slide, "CodeBuddy", 0.95, 6.05, 1.35, C_BLUE)
    pill(slide, "ClawPro", 2.45, 6.05, 1.25, C_TEAL)
    pill(slide, "WorkBuddy", 3.85, 6.05, 1.45, PptRGB(91, 84, 255))

    prs.save(PPTX)


def rebuild_zip():
    if STAGE.exists():
        shutil.rmtree(STAGE)
    STAGE.mkdir(parents=True)

    for dname in ["apps", "packages", "docs/submission", "docs/hackathons", "uipath", ".github/workflows"]:
        src = ROOT / dname
        if src.exists():
            dst = STAGE / dname
            dst.parent.mkdir(parents=True, exist_ok=True)
            shutil.copytree(src, dst, dirs_exist_ok=True)

    for fname in ["README.md", "LICENSE", "package.json", "package-lock.json", "tsconfig.base.json", "vitest.workspace.ts"]:
        src = ROOT / fname
        if src.exists():
            shutil.copy2(src, STAGE / fname)

    material = STAGE / "tencent-official-materials"
    material.mkdir(parents=True, exist_ok=True)
    shutil.copy2(DOCX, material / DOCX.name)
    shutil.copy2(PPTX, material / PPTX.name)
    shutil.copy2(ARCH, material / ARCH.name)
    (material / "README-Tencent-Submission.md").write_text(
        "# AgentGuard CI Tencent Submission Package\n\n"
        "This ZIP contains runnable source code, a polished Chinese whitepaper, a Tencent-themed pitch deck, and an architecture image.\n\n"
        "Recommended review path:\n\n"
        "1. Read `AgentGuard-CI-Tencent-Whitepaper.docx`.\n"
        "2. View `AgentGuard-CI-Tencent-Roadshow.pptx`.\n"
        "3. Run `npm install && npm test && npm run build`.\n"
        "4. Run `npm run agentguard:suite` and `npm run agentguard:agent-suite`.\n\n"
        "Public repository: https://github.com/baiqidi/agentguard-ci\n",
        encoding="utf-8",
    )

    if ZIP.exists():
        ZIP.unlink()
    with zipfile.ZipFile(ZIP, "w", zipfile.ZIP_DEFLATED) as archive:
        for path in STAGE.rglob("*"):
            if path.is_file():
                archive.write(path, path.relative_to(OUT))


if __name__ == "__main__":
    create_architecture_image()
    create_docx()
    create_pptx()
    rebuild_zip()
    print(json.dumps({
        "docx": str(DOCX),
        "docx_size": DOCX.stat().st_size,
        "pptx": str(PPTX),
        "pptx_size": PPTX.stat().st_size,
        "zip": str(ZIP),
        "zip_size": ZIP.stat().st_size,
        "architecture": str(ARCH),
    }, ensure_ascii=False, indent=2))
