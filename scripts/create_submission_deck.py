from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("docs/submission/AgentGuard-CI-deck.pptx")

NAVY = RGBColor(16, 31, 51)
TEAL = RGBColor(0, 145, 160)
MINT = RGBColor(118, 215, 196)
SLATE = RGBColor(67, 80, 98)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)


def add_title(slide, text, top=0.55, size=32):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(8.6), Inches(0.6))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.bold = True
    p.font.size = Pt(size)
    p.font.color.rgb = NAVY
    return box


def add_body(slide, bullets, left=0.75, top=1.55, width=8.3, size=18):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(4.7))
    tf = box.text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = SLATE
        p.space_after = Pt(9)
    return box


def add_badge(slide, text, left, top, width=2.1):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(0.42))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.color.rgb = TEAL
    tf = shape.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = WHITE
    return shape


def add_footer(slide):
    line = slide.shapes.add_shape(1, Inches(0), Inches(7.0), Inches(10), Inches(0.08))
    line.fill.solid()
    line.fill.fore_color.rgb = MINT
    line.line.color.rgb = MINT
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.08), Inches(8.8), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    p.text = "AgentGuard CI | UiPath AgentHack Track 3: UiPath Test Cloud"
    p.font.size = Pt(9)
    p.font.color.rgb = SLATE


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def title_slide(prs):
    slide = blank_slide(prs)
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.color.rgb = NAVY

    add_badge(slide, "UiPath Test Cloud", 0.72, 0.72, 2.25)
    title = slide.shapes.add_textbox(Inches(0.72), Inches(1.65), Inches(8.3), Inches(1.1))
    p = title.text_frame.paragraphs[0]
    p.text = "AgentGuard CI"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE

    subtitle = slide.shapes.add_textbox(Inches(0.76), Inches(2.75), Inches(8.3), Inches(1.0))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "Reliability testing for AI code-fixing agents"
    p.font.size = Pt(24)
    p.font.color.rgb = MINT

    body = slide.shapes.add_textbox(Inches(0.78), Inches(4.2), Inches(7.8), Inches(1.2))
    p = body.text_frame.paragraphs[0]
    p.text = "A governed scenario runner that turns agent fixes into auditable Test Cloud evidence."
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    return slide


def two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    slide = blank_slide(prs)
    add_title(slide, title)
    for x, heading, items in [(0.65, left_title, left_items), (5.12, right_title, right_items)]:
        panel = slide.shapes.add_shape(1, Inches(x), Inches(1.45), Inches(4.2), Inches(4.8))
        panel.fill.solid()
        panel.fill.fore_color.rgb = LIGHT
        panel.line.color.rgb = RGBColor(225, 231, 237)
        head = slide.shapes.add_textbox(Inches(x + 0.28), Inches(1.72), Inches(3.65), Inches(0.35))
        p = head.text_frame.paragraphs[0]
        p.text = heading
        p.font.bold = True
        p.font.size = Pt(17)
        p.font.color.rgb = NAVY
        add_body(slide, items, x + 0.35, 2.22, 3.45, 15)
    add_footer(slide)


def simple_slide(prs, title, bullets):
    slide = blank_slide(prs)
    add_title(slide, title)
    add_body(slide, bullets)
    add_footer(slide)


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs)

    two_column_slide(
        prs,
        "Problem",
        "AI agents can pass CI unsafely",
        [
            "A green build does not prove the agent made a safe fix.",
            "Agents may delete tests, edit unrelated files, or hide risky changes.",
            "Enterprise teams need repeatable evidence before approval.",
        ],
        "What judges should see",
        [
            "Clear reliability gates.",
            "Scenario-based evaluation.",
            "Artifacts that Test Cloud can store and review.",
        ],
    )

    simple_slide(
        prs,
        "Solution",
        [
            "AgentGuard CI runs controlled failure scenarios against a React and Express issue tracker.",
            "A code-fixing agent repairs the scenario, then AgentGuard scores the result across five gates.",
            "Each run emits Markdown, JSON, and JUnit evidence for governance and audit.",
        ],
    )

    two_column_slide(
        prs,
        "Architecture",
        "Prototype modules",
        [
            "apps/web: React issue tracker UI",
            "apps/api: Express API and triage logic",
            "packages/codefix-agent: scripted agent adapter",
            "packages/reliability-core: scoring and reports",
        ],
        "Evidence outputs",
        [
            "agentguard-runs/<scenario>/report.md",
            "agentguard-runs/<scenario>/result.json",
            "agentguard-runs/<scenario>/junit.xml",
            "uipath/test-cloud-import.csv",
        ],
    )

    simple_slide(
        prs,
        "Reliability Gates",
        [
            "Root cause quality: did the agent understand the failure?",
            "Test preservation: did it keep tests meaningful?",
            "Diff safety: did it avoid unrelated high-risk edits?",
            "CI health: do the required checks pass?",
            "Human approval readiness: is there enough evidence to approve?",
        ],
    )

    two_column_slide(
        prs,
        "Scenario Matrix",
        "Positive scenarios",
        [
            "frontend-contract: expected score 5/5",
            "backend-triage: expected score 5/5",
        ],
        "Governance scenarios",
        [
            "test-integrity-guard: fails by design when tests are weakened",
            "unsafe-diff-guard: fails by design on unrelated risky edits",
        ],
    )

    simple_slide(
        prs,
        "UiPath Test Cloud Fit",
        [
            "Each AgentGuard scenario maps to a Test Cloud test case.",
            "JUnit and Markdown reports become attached execution evidence.",
            "Governance failures are useful failures: they route unsafe agent behavior to human review.",
            "The result is a repeatable quality gate for AI agents, not a one-off demo.",
        ],
    )

    simple_slide(
        prs,
        "Demo Flow",
        [
            "1. Show the Issue Tracker sample app and CI failure.",
            "2. Run npm run agentguard:scenario -- --scenario frontend-contract.",
            "3. Show a 5/5 safe repair with scoped changes.",
            "4. Run unsafe-diff-guard and show governance failure evidence.",
            "5. Connect the reports to UiPath Test Cloud cases.",
        ],
    )

    simple_slide(
        prs,
        "Next Steps",
        [
            "Connect the runner to UiPath Orchestrator and Test Cloud APIs.",
            "Add adapters for more coding agents and enterprise CI providers.",
            "Expand the scenario library for security, data integrity, and compliance workflows.",
            "Use the same evidence loop for production agent release gates.",
        ],
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
